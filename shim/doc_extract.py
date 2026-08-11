"""EDM 文档文本提取（issue #48，自研）：bytes-in / text-out，服务端不落原始文件盘。

上传的二进制文档（PDF/DOCX/XLSX/PPTX）提取为纯文本后交 admin_api 走 edm_lib 指纹管线；
文本类（.txt/.md 等）按多编码链解码（GBK 中文文档可正确解码，护栏不回退）。
提取结果不含任何合成标记行（如 "--- Page N ---"）——指纹按真实内容计算，
避免标记行在跨文档间产生相同的行级指纹（误命中）或让空文档混过入库最小长度门槛。
扫描版 PDF（内嵌文本为空）与图片走 Tesseract OCR（issue #50，本地识别不出域）：
图片直接 `pytesseract.image_to_string(lang="chi_sim+eng")`；扫描 PDF 用 PyMuPDF 渲染页图
（180 DPI）逐页 OCR，页数上限 MAX_OCR_PAGES。tesseract 缺失/执行失败 → OcrUnavailableError；
OCR 全空 → EmptyDocumentError。OCR 质量边界：中文印刷体一般、手写差、表格版面丢失。
水印文本层回退（issue #52 缺口 2，真实试点「扫描全能王 创建」式水印）：内嵌文本非空但
有效字符/页 < OCR_FALLBACK_MIN_CHARS_PER_PAGE 且页面含扫描图（issue #53 P1-1 按实际显示
面积判定：单图 bbox > 页面面积 SCAN_IMAGE_AREA_RATIO 才算扫描页，logo/签名章小图不触发）
时仍回退 OCR——水印层非空会挡住朴素 OCR 回退，正文指纹全丢。docx/pptx 提取覆盖表格与文本框（issue #52 缺口 1）：
python-docx iter_inner_content() 按文档顺序遍历段落+表格（旧实现 document.paragraphs
丢失全部表格文本，试点真实文档 95% 内容在表格）；pptx 递归组合 shape、表格按行拼单元格。
OCR 资源护栏（issue #51 评审修复）：渲染/解码目标像素两路径共用 MAX_OCR_PIXELS 上限
（PDF 页按 page.rect×DPI 系数渲染前拦，直传图片按 Image.open 读到的 size 解码前拦）——
数万 pt MediaBox 的畸形 PDF 仅几 KB，三重字节/页数/字符上限都兜不住其渲染像素量，
超限可 OOM 打挂 shim（检测链 fail-open 窗口，#49 P1-1 同类）；tesseract 单次执行带
OCR_PAGE_TIMEOUT 超时（pytesseract timeout kwarg），超时/引擎失败 → OcrUnavailableError；
直传图片路径 OCR 调用兜底一切内部异常转 CorruptDocumentError（不裸泄断线）。
提取文本长度上限 MAX_EXTRACTED_CHARS（issue #49 P1-1）：zip 压缩态/流扩张可使 16MB 文件
提出数十 MB 文本，edm_lib.shingles 全量滑窗内存随字符数线性膨胀，无上限可 OOM 打挂 shim
（检测链 fail-open 窗口）；8M 字符 ≈ 500 页 PDF（~1.5M 字符）的 5 倍余量。
第三方解析库（fitz/docx/openpyxl/pptx/pytesseract/Pillow）一律函数级懒加载（issue #49 P2-7）：
app.py → admin_api → doc_extract 模块级 import 链不含第三方库，解析库缺失/损坏不波及
/request /response 检测路径。
"""
import io
from pathlib import PurePosixPath

OFFICE_EXTENSIONS = frozenset({".pdf", ".docx", ".xlsx", ".pptx"})
TEXT_EXTENSIONS = frozenset({".txt", ".text", ".log", ".md", ".markdown"})
# OCR 图片（issue #50）：本地 Tesseract 识别，chi_sim+eng
IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif"})
SUPPORTED_EXTENSIONS = OFFICE_EXTENSIONS | TEXT_EXTENSIONS | IMAGE_EXTENSIONS
_LEGACY_DOC_EXTENSIONS = frozenset({".doc"})
_UNSUPPORTED_IMAGE_EXTENSIONS = frozenset({".gif", ".webp"})  # Pillow 可读但本期不纳管，明确报错

# utf-8-sig 在前剥掉 BOM；GBK 系兜底中文 Windows 文档；latin-1 全字节可解（最终兜底，不抛错）
_TEXT_DECODINGS = ("utf-8-sig", "utf-8", "gbk", "gb18030", "latin-1")

_PDF_MAGIC = b"%PDF-"
_OOXML_MAGIC = b"PK\x03\x04"  # docx/xlsx/pptx 均为 zip 包
_IMAGE_MAGICS = (  # 与 IMAGE_EXTENSIONS 对应：png / jpeg / bmp / tiff(LE|BE)
    b"\x89PNG\r\n\x1a\n",
    b"\xff\xd8\xff",
    b"BM",
    b"II*\x00",
    b"MM\x00*",
)

# 提取文本字符数上限（issue #49 P1-1）：防 zip/流扩张提出超大文本 OOM 打挂 shim
MAX_EXTRACTED_CHARS = 8 * 1000 * 1000
MAX_OCR_PAGES = 50  # 扫描 PDF OCR 页数上限（issue #50）：单页 OCR 秒级，防超长扫描件拖死上传请求
# OCR 单页/单图渲染-解码像素预算（issue #51 P1-1，两条 OCR 路径共用）：≈180DPI 的 A0 页有余量；
# 数万 pt MediaBox 畸形 PDF（文件仅几 KB，字节/页数/字符三重上限兜不住）渲染即数百亿像素，
# 直传超大图解码同理，超限可 OOM 打挂 shim（检测链 fail-open 窗口，#49 P1-1 同类）
MAX_OCR_PIXELS = 40_000_000
# tesseract 单次执行超时秒数（issue #51 P2-2，pytesseract timeout kwarg，单页/单图同值）：
# 防引擎挂死拖住管理端会话、线程占用无界
OCR_PAGE_TIMEOUT = 60
# 水印文本层密度阈值（issue #52 缺口 2）：内嵌有效字符/页低于此值且页面含扫描图 → 回退 OCR。
# 「扫描全能王 创建」式水印每页仅数字符（试点真实合同 8 字符/页）；正常数字 PDF 每页数百字符起
OCR_FALLBACK_MIN_CHARS_PER_PAGE = 50
# 扫描页判定面积比（issue #53 P1-1）：单图实际显示 bbox 面积 > 页面面积×此值才算扫描页。
# get_images() 只看 /Resources 引用——信头 logo/签名章小图都算「含图」，稀疏数字单页
# （<50 字符/页）+ logo 会误触发回退：内嵌完美文本被换成 OCR 次品（指纹质量静默退化），
# OCR 全空或引擎缺失时甚至 400 拒收（#52 前可正常入库，行为回归）。get_image_info()
# 只含实际显示的图片，按面积过滤——试点扫描件整页图实测 ≈86%，logo/签名章 <1%
SCAN_IMAGE_AREA_RATIO = 0.5
_OCR_DPI = 180
_OCR_LANG = "chi_sim+eng"


class DocumentExtractionError(Exception):
    """提取失败基类：str(exc) 为用户可见中文文案。"""


class UnsupportedDocumentError(DocumentExtractionError):
    """扩展名不支持（含 .doc 老式格式与图片的定制提示）。"""


class CorruptDocumentError(DocumentExtractionError):
    """内容与扩展名不符、加密或解析失败（含图片 OCR 路径内部异常兜底，issue #51 P2-3）。"""


class EmptyDocumentError(DocumentExtractionError):
    """空文件或未提取到文本（含 OCR 全空——手写/低质量图识别率差）。"""


class OcrUnavailableError(DocumentExtractionError):
    """OCR 引擎不可用（issue #50）：tesseract 二进制缺失、执行失败或超时（issue #51 P2-2）。"""


class OcrImageTooLargeError(DocumentExtractionError):
    """OCR 渲染/解码目标像素超过 MAX_OCR_PIXELS（issue #51 P1-1）。"""


class ExtractedTextTooLargeError(DocumentExtractionError):
    """提取文本超过 MAX_EXTRACTED_CHARS（issue #49 P1-1）。"""


def extract_text_from_bytes(filename: str, data: bytes) -> str:
    """从单个文档的原始字节提取纯文本；失败抛 DocumentExtractionError 子类。"""
    if not data:
        raise EmptyDocumentError(f"「{filename}」是空文件")
    ext = PurePosixPath(filename or "").suffix.lower()
    if ext in _LEGACY_DOC_EXTENSIONS:
        raise UnsupportedDocumentError(
            "老式 .doc（二进制格式）暂不支持：请在 Word/WPS 中另存为 .docx 后上传"
        )
    if ext in _UNSUPPORTED_IMAGE_EXTENSIONS:
        raise UnsupportedDocumentError("GIF/WebP 图片暂不支持：请转 PNG/JPG 后上传")
    if ext not in SUPPORTED_EXTENSIONS:
        raise UnsupportedDocumentError(
            f"不支持的文件类型 '{ext or filename}'：支持 .pdf/.docx/.xlsx/.pptx、"
            ".txt/.md/.text/.log 与 .png/.jpg/.jpeg/.bmp/.tiff/.tif 图片（OCR）"
        )

    if ext == ".pdf":
        text = _extract_pdf(data, filename)
    elif ext == ".docx":
        text = _extract_docx(data, filename)
    elif ext == ".xlsx":
        text = _extract_xlsx(data, filename)
    elif ext == ".pptx":
        text = _extract_pptx(data, filename)
    elif ext in IMAGE_EXTENSIONS:
        text = _extract_image_ocr(data, filename)
    else:  # 文本类
        text = _decode_text(data)

    if not text.strip():
        # 空结果文案分路径（issue #51 P2-4）：OCR 路径（图片/扫描 PDF）给 OCR 质量提示，
        # 非 OCR 空文档（.docx/.txt 等）给原文案——PDF 文本层非空已在 _extract_pdf 提前返回，
        # 走到这里的 .pdf 必是 OCR 路径
        if ext in IMAGE_EXTENSIONS or ext == ".pdf":
            raise EmptyDocumentError(
                f"「{filename}」未提取到文本（扫描件/图片 OCR 对手写体与低清晰度识别率差，"
                "请用更清晰版本或文字版文件）"
            )
        raise EmptyDocumentError(f"「{filename}」未提取到文本")
    if len(text) > MAX_EXTRACTED_CHARS:
        raise ExtractedTextTooLargeError(
            f"「{filename}」提取文本 {len(text)} 字符超过 {MAX_EXTRACTED_CHARS // 1000000}00 万字符上限"
            "（zip/流扩张防 OOM 保护），请拆分文档后上传"
        )
    return text


def _check_magic(data: bytes, magic: bytes, filename: str, kind: str) -> None:
    """文件头魔数校验：拦扩展名伪造（把 zip 改名为 .pdf 之类）。"""
    if not data.startswith(magic):
        raise CorruptDocumentError(f"「{filename}」内容不是有效 {kind}（文件头与扩展名不符）")


def _extract_pdf(data: bytes, filename: str) -> str:
    import fitz  # PyMuPDF；懒加载（issue #49 P2-7）：解析库缺失不波及模块 import

    _check_magic(data, _PDF_MAGIC, filename, "PDF")
    try:
        with fitz.open(stream=data, filetype="pdf") as doc:
            if doc.is_encrypted and not doc.authenticate(""):
                raise CorruptDocumentError(f"「{filename}」已加密，无法读取")
            text = "\n\n".join(page.get_text() for page in doc)
            if text.strip() and not _needs_ocr_fallback(doc, text):
                return text
            # 无内嵌文本层（issue #50），或水印文本层挡住（issue #52 缺口 2）→ 渲染页图逐页 OCR
            return _ocr_pdf_pages(doc, filename)
    except (CorruptDocumentError, OcrUnavailableError, OcrImageTooLargeError, EmptyDocumentError):
        raise
    except Exception as e:
        raise CorruptDocumentError(f"「{filename}」PDF 解析失败（{e}）") from e


def _needs_ocr_fallback(doc, text: str) -> bool:
    """水印文本层判定（issue #52 缺口 2）：有效字符/页 < OCR_FALLBACK_MIN_CHARS_PER_PAGE
    且页面含扫描图 → 回退 OCR。「扫描全能王 创建」式水印页页几字符，非空文本层会挡住
    朴素 OCR 回退，正文指纹全丢；无扫描图的稀疏文本页（数字文档留白多）不回退——OCR 无可识别对象。"""
    effective = len("".join(text.split()))  # 有效字符：剔除全部空白
    if effective / max(doc.page_count, 1) >= OCR_FALLBACK_MIN_CHARS_PER_PAGE:
        return False
    return any(_is_scan_page(page) for page in doc)


def _is_scan_page(page) -> bool:
    """单图实际显示 bbox 面积 > 页面面积×SCAN_IMAGE_AREA_RATIO 才算扫描页（issue #53 P1-1）：
    get_image_info() 仅含实际显示的图片，logo/签名章小图（实测 <1%）不触发回退；
    get_images() 看 /Resources 引用不过滤显示尺寸，会误伤稀疏数字单页。"""
    page_area = page.rect.width * page.rect.height
    if page_area <= 0:
        return False
    for info in page.get_image_info():
        x0, y0, x1, y1 = info["bbox"]
        if (x1 - x0) * (y1 - y0) > SCAN_IMAGE_AREA_RATIO * page_area:
            return True
    return False


def _ocr_pdf_pages(doc, filename: str) -> str:
    """扫描 PDF 逐页 OCR：PyMuPDF 渲染 180 DPI 页图 → Tesseract。页数/像素超上限明确报错。"""
    if doc.page_count > MAX_OCR_PAGES:
        raise EmptyDocumentError(
            f"「{filename}」无有效文本层需 OCR（扫描件或仅水印文本），共 {doc.page_count} 页超过 "
            f"{MAX_OCR_PAGES} 页上限（防超长扫描件拖死上传），请拆分文档后上传"
        )
    from PIL import Image  # 懒加载（issue #49 P2-7 纪律覆盖 OCR 依赖）

    zoom = _OCR_DPI / 72
    import fitz  # 已在 _extract_pdf 加载；局部再引保持函数自含
    pages = []
    for index, page in enumerate(doc, start=1):
        # 像素预算（issue #51 P1-1）：渲染前按 page.rect×DPI 系数算目标像素数，
        # 超限先拦——畸形大 MediaBox 页渲染即 OOM，等 get_pixmap 抛错就晚了
        pixels = page.rect.width * zoom * page.rect.height * zoom
        if pixels > MAX_OCR_PIXELS:
            raise OcrImageTooLargeError(
                f"「{filename}」第 {index} 页版面过大：180DPI 渲染需约 {pixels / 1000000:.0f}00 万像素，"
                f"超过 {MAX_OCR_PIXELS // 1000000}00 万像素上限（防超大版面渲染 OOM 保护），"
                "请降低扫描分辨率或拆分文档后上传"
            )
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        pages.append(_run_tesseract(img, filename))
    return "\n\n".join(pages)


def _extract_image_ocr(data: bytes, filename: str) -> str:
    """图片 OCR（issue #50）：魔数校验 → Pillow 打开 → Tesseract chi_sim+eng。
    像素上限在解码前校验（Image.open 只读头即得 size，issue #51 P1-1）；
    OCR 调用兜底一切内部异常转中文报错（issue #51 P2-3，对齐 PDF 路径「提取失败不裸泄」惯例）。"""
    if not any(data.startswith(m) for m in _IMAGE_MAGICS):
        raise CorruptDocumentError(f"「{filename}」内容不是有效图片（文件头与扩展名不符）")
    from PIL import Image  # 懒加载

    try:
        img = Image.open(io.BytesIO(data))
        if img.width * img.height > MAX_OCR_PIXELS:
            raise OcrImageTooLargeError(
                f"「{filename}」图片尺寸 {img.width}×{img.height} 超过 "
                f"{MAX_OCR_PIXELS // 1000000}00 万像素上限（防超大图解码 OOM 保护），请压缩后上传"
            )
        img.load()
    except DocumentExtractionError:
        raise
    except Exception as e:
        raise CorruptDocumentError(f"「{filename}」图片打开失败（{e}）") from e
    try:
        return _run_tesseract(img, filename)
    except DocumentExtractionError:
        raise
    except Exception as e:  # pytesseract/Pillow 内部异常（OSError/MemoryError/TypeError 等）
        raise CorruptDocumentError(f"「{filename}」图片 OCR 识别失败（{e}）") from e


def _run_tesseract(img, filename: str) -> str:
    """单图 Tesseract 识别；引擎缺失/执行失败/超时 → OcrUnavailableError（明确中文报错）。"""
    import pytesseract  # 懒加载

    try:
        return pytesseract.image_to_string(img, lang=_OCR_LANG, timeout=OCR_PAGE_TIMEOUT)
    except pytesseract.TesseractNotFoundError as e:
        raise OcrUnavailableError(
            "OCR 引擎（tesseract）不可用：容器未安装或不在 PATH，无法处理扫描件/图片"
        ) from e
    except pytesseract.TesseractError as e:
        raise OcrUnavailableError(f"「{filename}」OCR 执行失败（{e}）") from e
    except RuntimeError as e:  # pytesseract timeout kwarg 超时即抛 RuntimeError（issue #51 P2-2）
        raise OcrUnavailableError(
            f"「{filename}」OCR 识别超时（单页/单图超过 {OCR_PAGE_TIMEOUT} 秒上限），"
            "请降低图片清晰度或页面复杂度后重试"
        ) from e


def _extract_docx(data: bytes, filename: str) -> str:
    from docx import Document as DocxDocument  # 懒加载（issue #49 P2-7）
    from docx.table import Table as DocxTable

    _check_magic(data, _OOXML_MAGIC, filename, "Office 文件")
    try:
        doc = DocxDocument(io.BytesIO(data))
        # iter_inner_content()（python-docx 1.2.0）按文档顺序产出段落+表格（issue #52 缺口 1：
        # 旧实现只取 document.paragraphs，试点真实文档 95% 内容在表格里全丢）
        lines = []
        for block in doc.iter_inner_content():
            if isinstance(block, DocxTable):
                for row in block.rows:
                    line = "\t".join(_dedup_adjacent_cells([cell.text.strip() for cell in row.cells]))
                    if line.strip():
                        lines.append(line)
            elif block.text.strip():
                lines.append(block.text)
        return "\n".join(lines)
    except Exception as e:
        raise CorruptDocumentError(f"「{filename}」DOCX 解析失败（{e}）") from e


def _dedup_adjacent_cells(cells: list) -> list:
    """行内相邻重复 cell 文本去重（保序，issue #53 P2-1）：python-docx/pptx 水平合并单元格
    row.cells 会重复返回同一合并 cell（["A","A","C"]），原样提取产生 "A\tA\tC" 式重复段，
    与员工粘贴原文的 shingle 窗口失配。跨行重复（垂直合并产物）保守不动——相邻整行重复
    也可能是真实重复数据行，误去重会丢正文。"""
    out = []
    for cell in cells:
        if not out or out[-1] != cell:
            out.append(cell)
    return out


def _extract_xlsx(data: bytes, filename: str) -> str:
    from openpyxl import load_workbook  # 懒加载（issue #49 P2-7）

    _check_magic(data, _OOXML_MAGIC, filename, "Office 文件")
    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        try:
            lines = []
            for sheet_name in wb.sheetnames:
                for row in wb[sheet_name].iter_rows(values_only=True):
                    line = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if line.strip():
                        lines.append(line)
            return "\n".join(lines)
        finally:
            wb.close()
    except Exception as e:
        raise CorruptDocumentError(f"「{filename}」XLSX 解析失败（{e}）") from e


def _extract_pptx(data: bytes, filename: str) -> str:
    from pptx import Presentation as PptxPresentation  # 懒加载（issue #49 P2-7）

    _check_magic(data, _OOXML_MAGIC, filename, "Office 文件")
    try:
        prs = PptxPresentation(io.BytesIO(data))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                lines.extend(_pptx_shape_lines(shape))
        return "\n".join(lines)
    except Exception as e:
        raise CorruptDocumentError(f"「{filename}」PPTX 解析失败（{e}）") from e


def _pptx_shape_lines(shape) -> list:
    """单 shape 的文本行（issue #52 缺口 1 补齐）：组合 shape 递归；表格按行拼单元格；
    文本框/占位符取 shape.text。旧实现 getattr(shape,'text','') 对表格/组合一律落空。"""
    if hasattr(shape, "shapes"):  # 组合 shape：递归取子 shape
        lines = []
        for sub in shape.shapes:
            lines.extend(_pptx_shape_lines(sub))
        return lines
    if getattr(shape, "has_table", False):  # 表格（GraphicFrame）
        lines = []
        for row in shape.table.rows:
            line = "\t".join(_dedup_adjacent_cells([cell.text.strip() for cell in row.cells]))
            if line.strip():
                lines.append(line)
        return lines
    text = getattr(shape, "text", "")
    return [text] if text.strip() else []


def _decode_text(data: bytes) -> str:
    """文本类多编码链解码（GBK 护栏）：逐个尝试，latin-1 兜底保证不抛 UnicodeDecodeError。"""
    for encoding in _TEXT_DECODINGS:
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")  # pragma: no cover - latin-1 不失败
