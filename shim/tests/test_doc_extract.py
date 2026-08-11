#!/usr/bin/env python3
"""shim/doc_extract.py（issue #48，自研提取器）单测。

覆盖：四种 Office 格式提取（含中文）、文本类多编码解码（GBK/BOM）、
拒绝路径（.doc/图片/未知扩展名/空文件/魔数不符/损坏 zip/扫描 PDF）。
fixture 全部现场生成（python-docx/openpyxl/python-pptx/PyMuPDF），不落盘；
造档函数放模块级，admin API 端到端用例（test_admin_api.py）复用。
"""
import io
import os
import sys
import unittest

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import doc_extract as dx  # noqa: E402


# ---------------------------------------------------------------------------
# fixture：现场生成各格式文档字节
# ---------------------------------------------------------------------------


def make_docx_bytes(paragraphs) -> bytes:
    from docx import Document

    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_xlsx_bytes(sheets) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    first = True
    for name, rows in sheets.items():
        ws = wb.active if first else wb.create_sheet()
        ws.title = name
        for row in rows:
            ws.append(row)
        first = False
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_pptx_bytes(slides_text) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for slide_texts in slides_text:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 近空白版式
        for i, text in enumerate(slide_texts):
            tb = slide.shapes.add_textbox(Inches(1), Inches(1 + i * 0.6), Inches(6), Inches(0.5))
            tb.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_pdf_bytes(pages_text, fontname="helv") -> bytes:
    """一页一段文字；fontname='china-s' 用 PyMuPDF 内置 CJK 字体写中文。"""
    import fitz

    doc = fitz.open()
    for text in pages_text:
        page = doc.new_page()
        if text:
            page.insert_text((72, 72), text, fontname=fontname)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def make_png_bytes(text, fontname="helv") -> bytes:
    """fitz 渲染文字为 PNG 图片字节（无文本层，放大 2 倍利于 OCR）；
    fontname='china-s' 可渲染中文（不依赖宿主 CJK 字体）。页面 3200pt 宽，容 ~60 个中文字符。"""
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=3200, height=300)
    page.insert_text((60, 150), text, fontname=fontname, fontsize=48)
    data = page.get_pixmap(matrix=fitz.Matrix(2, 2)).tobytes("png")
    doc.close()
    return data


def make_scanned_pdf_bytes(pages_text, fontname="helv") -> bytes:
    """文字渲染成图后存入新 PDF（无内嵌文本层），模拟扫描件。fontsize 16 保证长行不溢出页宽。"""
    import fitz

    src = fitz.open()
    for t in pages_text:
        page = src.new_page()
        if t:
            page.insert_text((72, 72), t, fontname=fontname, fontsize=16)
    out = fitz.open()
    for page in src:
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        new_page = out.new_page(width=page.rect.width, height=page.rect.height)
        new_page.insert_image(new_page.rect, pixmap=pix)
    buf = io.BytesIO()
    out.save(buf)
    out.close()
    src.close()
    return buf.getvalue()


def make_watermarked_scanned_pdf_bytes(pages_text, watermark="扫描全能王 创建") -> bytes:
    """扫描图 + 真文本水印同页（issue #52 缺口 2 样本，「扫描全能王」式）：正文只以页图存在，
    内嵌文本层每页仅水印数字符——非空文本层会挡住朴素 OCR 回退，正文指纹全丢。"""
    import fitz

    doc = fitz.open()
    for t in pages_text:
        src = fitz.open()
        sp = src.new_page()
        sp.insert_text((72, 72), t, fontsize=16)
        pix = sp.get_pixmap(matrix=fitz.Matrix(2, 2))
        page = doc.new_page(width=sp.rect.width, height=sp.rect.height)
        page.insert_image(page.rect, pixmap=pix)  # 正文仅以扫描图存在
        page.insert_text((72, 750), watermark, fontname="china-s", fontsize=9)  # 真文本水印层
        src.close()
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def make_pdf_with_image_bytes(text, fontname="helv") -> bytes:
    """密集数字文本 + 内嵌小图（issue #52 不回退分支样本）：密度 ≥ 阈值，即使有图也不走 OCR。
    中文文本需 fontname='china-s'（helv 渲染中文变 '·'）。"""
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    # insert_textbox 在矩形内自动换行（单行超页宽会被 fitz 截断，密度样本必须用盒式写入）
    page.insert_textbox(fitz.Rect(72, 72, 540, 400), text, fontname=fontname, fontsize=12)
    img = fitz.open()
    ip = img.new_page(width=100, height=100)
    pix = ip.get_pixmap()
    page.insert_image(fitz.Rect(72, 400, 172, 500), pixmap=pix)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    img.close()
    return buf.getvalue()


def make_docx_with_table_bytes(before, table_rows, after) -> bytes:
    """段落-表格-段落交错（issue #52 缺口 1 样本）：检验 iter_inner_content 文档顺序与表格提取。"""
    from docx import Document

    doc = Document()
    doc.add_paragraph(before)
    table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
    for i, row in enumerate(table_rows):
        for j, cell_text in enumerate(row):
            table.cell(i, j).text = cell_text
    doc.add_paragraph(after)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_docx_merged_table_bytes() -> bytes:
    """水平合并单元格 docx（issue #53 P2-1 样本）：合并后 row.cells 重复返回同一 cell。"""
    from docx import Document

    doc = Document()
    table = doc.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "合并甲"
    table.cell(0, 0).merge(table.cell(0, 1))  # 水平合并 → row.cells = [c, c, 丙]
    table.cell(0, 2).text = "丙"
    table.cell(1, 0).text = "丁"
    table.cell(1, 1).text = "戊"
    table.cell(1, 2).text = "己"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_mixed_density_pdf_bytes(pages) -> bytes:
    """多页混合密度样本（issue #53 P2-3 边界用例）：pages = [(text, full_page_image), ...]。
    full_page_image=True 时先铺整页图（扫描页，面积比 100%）；文本以换行盒式写入文本层。"""
    import fitz

    doc = fitz.open()
    for text, full_page_image in pages:
        page = doc.new_page()
        if full_page_image:
            img = fitz.open()
            ip = img.new_page(width=200, height=200)
            page.insert_image(page.rect, pixmap=ip.get_pixmap())
            img.close()
        if text:
            page.insert_textbox(fitz.Rect(72, 72, 540, 760), text, fontname="china-s", fontsize=12)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def make_pptx_with_table_bytes(textbox_text, table_rows) -> bytes:
    """文本框 + 表格同页（issue #52 pptx 缺口样本）：旧实现对表格 GraphicFrame 一律取空文本。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(0.5))
    tb.text_frame.text = textbox_text
    graphic = slide.shapes.add_table(len(table_rows), len(table_rows[0]),
                                     Inches(1), Inches(2), Inches(6), Inches(1.5))
    for i, row in enumerate(table_rows):
        for j, cell_text in enumerate(row):
            graphic.table.cell(i, j).text = cell_text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_huge_mediabox_pdf_bytes(side_pt=20000) -> bytes:
    """单页空白 PDF（无文本层），MediaBox 改成 side_pt 见方（issue #51 P1-1 畸形大版面样本）：
    文件仅几百字节，180DPI 渲染却需 (side_pt×2.5)² 像素（20000pt → 25 亿 px）。"""
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=612, height=792)
    page.set_mediabox(fitz.Rect(0, 0, side_pt, side_pt))
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def make_png_with_dimensions(width, height) -> bytes:
    """最小 PNG：IHDR 声明 width×height，IDAT 仅 1 像素数据（不打算被解码）。
    用于像素上限用例（issue #51 P1-1）：Pillow Image.open 只读头即得 size，
    服务端须在 img.load() 解码前拦截，因此伪造尺寸无需真实大图。"""
    import struct
    import zlib

    def chunk(tag, payload):
        return (struct.pack(">I", len(payload)) + tag + payload
                + struct.pack(">I", zlib.crc32(tag + payload)))

    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)  # 8-bit RGB
    idat = zlib.compress(b"\x00" + b"\x00\x00\x00")  # 1 scanline(filter 0) + 1 黑像素
    return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr) + chunk(b"IDAT", idat) + chunk(b"IEND", b"")


def _tesseract_available() -> bool:
    import shutil

    return shutil.which("tesseract") is not None


def _chi_sim_available() -> bool:
    import shutil
    import subprocess

    if not shutil.which("tesseract"):
        return False
    r = subprocess.run(["tesseract", "--list-langs"], capture_output=True, text=True)
    return "chi_sim" in (r.stdout + r.stderr)


# ---------------------------------------------------------------------------
# 各格式提取
# ---------------------------------------------------------------------------


class TestExtractDocx(unittest.TestCase):
    def test_basic(self):
        text = dx.extract_text_from_bytes(
            "doc.docx", make_docx_bytes(["Hello world", "Second paragraph", "  "])
        )
        self.assertIn("Hello world", text)
        self.assertIn("Second paragraph", text)

    def test_chinese(self):
        text = dx.extract_text_from_bytes("合同.docx", make_docx_bytes(["机密合同条款 第一条"]))
        self.assertIn("机密合同条款 第一条", text)

    def test_table_in_document_order(self):
        """docx 表格提取（issue #52 缺口 1）：iter_inner_content 按文档顺序遍历段落+表格，
        表格按行拼单元格——试点真实文档 95% 内容在表格，旧实现（仅 document.paragraphs）全丢。"""
        data = make_docx_with_table_bytes(
            "合同首部段落",
            [["甲方", "甲公司"], ["乙方", "乙公司"]],
            "合同尾部段落",
        )
        text = dx.extract_text_from_bytes("contract.docx", data)
        self.assertIn("甲方\t甲公司", text)
        self.assertIn("乙方\t乙公司", text)
        # 文档顺序：首部段落 < 表格行 < 尾部段落
        self.assertLess(text.index("合同首部段落"), text.index("甲方"))
        self.assertLess(text.index("乙方\t乙公司"), text.index("合同尾部段落"))

    def test_merged_cells_dedup(self):
        """docx 合并单元格相邻去重（issue #53 P2-1）：水平合并的重复 cell 保序去重，
        不产生「合并甲\t合并甲\t丙」式重复段（与员工粘贴原文的 shingle 窗口失配）。"""
        text = dx.extract_text_from_bytes("merged.docx", make_docx_merged_table_bytes())
        self.assertIn("合并甲\t丙", text)
        self.assertNotIn("合并甲\t合并甲", text)
        self.assertIn("丁\t戊\t己", text)  # 未合并行原样


class TestExtractXlsx(unittest.TestCase):
    def test_multiple_sheets(self):
        data = make_xlsx_bytes({"Alpha": [["a1", "b1"], ["a2", 42]], "Beta": [["x", "y"]]})
        text = dx.extract_text_from_bytes("book.xlsx", data)
        self.assertIn("a1", text)
        self.assertIn("42", text)
        self.assertIn("x", text)

    def test_chinese(self):
        text = dx.extract_text_from_bytes("报表.xlsx", make_xlsx_bytes({"数据": [["项目", "金额"], ["阿尔法", 42]]}))
        self.assertIn("项目", text)
        self.assertIn("阿尔法", text)


class TestExtractPptx(unittest.TestCase):
    def test_basic(self):
        data = make_pptx_bytes([["Slide 1 title", "Slide 1 body"], ["Slide 2 only text"]])
        text = dx.extract_text_from_bytes("deck.pptx", data)
        self.assertIn("Slide 1 title", text)
        self.assertIn("Slide 2 only text", text)

    def test_chinese(self):
        text = dx.extract_text_from_bytes("路演.pptx", make_pptx_bytes([["路演机密数据 2026"]]))
        self.assertIn("路演机密数据 2026", text)

    def test_textbox_and_table(self):
        """pptx 文本框+表格提取（issue #52 缺口 1 核查补齐）：旧实现 getattr(shape,'text','')
        对表格 GraphicFrame 一律落空。"""
        data = make_pptx_with_table_bytes("季度汇报标题", [["指标", "数值"], ["违约金", "百分之二十"]])
        text = dx.extract_text_from_bytes("deck.pptx", data)
        self.assertIn("季度汇报标题", text)
        self.assertIn("指标\t数值", text)
        self.assertIn("违约金\t百分之二十", text)


class TestExtractPdf(unittest.TestCase):
    def test_basic(self):
        text = dx.extract_text_from_bytes("a.pdf", make_pdf_bytes(["Hello PDF world"]))
        self.assertIn("Hello PDF world", text)

    def test_chinese(self):
        text = dx.extract_text_from_bytes("z.pdf", make_pdf_bytes(["机密PDF内容"], fontname="china-s"))
        self.assertIn("机密PDF内容", text)

    def test_multipage(self):
        text = dx.extract_text_from_bytes("m.pdf", make_pdf_bytes(["page one text", "page two text"]))
        self.assertIn("page one text", text)
        self.assertIn("page two text", text)

    def test_no_synthetic_page_markers(self):
        """提取结果不含 "--- Page N ---" 类合成标记（指纹按真实内容计算，防跨文档误命中）。"""
        text = dx.extract_text_from_bytes("m.pdf", make_pdf_bytes(["page one text", "page two text"]))
        self.assertNotIn("---", text)


class TestExtractText(unittest.TestCase):
    def test_utf8(self):
        text = dx.extract_text_from_bytes("note.txt", "hello 你好\nline two".encode("utf-8"))
        self.assertIn("hello 你好", text)
        self.assertIn("line two", text)

    def test_gbk(self):
        """GBK 护栏：中文 Windows 文档正确解码（GBK 字节现场编码，源码不贴不可见字符）。"""
        text = dx.extract_text_from_bytes("note.txt", "你好世界".encode("gbk"))
        self.assertEqual(text, "你好世界")

    def test_utf8_bom_stripped(self):
        text = dx.extract_text_from_bytes("note.md", b"\xef\xbb\xbfhello")
        self.assertEqual(text, "hello")

    def test_markdown_ext(self):
        self.assertIn("# T", dx.extract_text_from_bytes("d.markdown", b"# T\n\nbody\n"))


# ---------------------------------------------------------------------------
# 拒绝路径
# ---------------------------------------------------------------------------


class TestRejections(unittest.TestCase):
    def test_empty_bytes(self):
        with self.assertRaises(dx.EmptyDocumentError):
            dx.extract_text_from_bytes("f.docx", b"")

    def test_legacy_doc(self):
        with self.assertRaises(dx.UnsupportedDocumentError) as cm:
            dx.extract_text_from_bytes("old.doc", b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1 ole")
        self.assertIn(".docx", str(cm.exception))

    def test_gif_webp_image(self):
        """GIF/WebP 不纳入 OCR 面（issue #50 范围 png/jpg/jpeg/bmp/tiff），明确报错转 PNG/JPG。"""
        for fn in ("a.gif", "a.webp"):
            with self.subTest(fn=fn), self.assertRaises(dx.UnsupportedDocumentError) as cm:
                dx.extract_text_from_bytes(fn, b"GIF89a" + b"\x00" * 32)
            self.assertIn("PNG/JPG", str(cm.exception))

    def test_image_magic_mismatch(self):
        """图片魔数校验：内容不是有效图片 → Corrupt（无需 tesseract）。"""
        with self.assertRaises(dx.CorruptDocumentError):
            dx.extract_text_from_bytes("pic.png", b"this is not a png at all.........")

    def test_unknown_extension(self):
        with self.assertRaises(dx.UnsupportedDocumentError):
            dx.extract_text_from_bytes("a.zip", b"PK\x03\x04" + b"\x00" * 32)

    def test_pdf_magic_mismatch(self):
        with self.assertRaises(dx.CorruptDocumentError):
            dx.extract_text_from_bytes("f.pdf", b"this is not a pdf")

    def test_ooxml_magic_mismatch(self):
        with self.assertRaises(dx.CorruptDocumentError):
            dx.extract_text_from_bytes("f.docx", b"not an office file")

    def test_corrupt_docx_zip(self):
        """魔数对头（PK）但 zip 体损坏 → Corrupt。"""
        with self.assertRaises(dx.CorruptDocumentError):
            dx.extract_text_from_bytes("f.docx", b"PK\x03\x04" + b"\x00" * 512)

    def test_scanned_pdf_blank_page_ocr_path(self):
        """扫描版 PDF（空白页无文本层）走 OCR 路径（issue #50）：引擎缺失 → OcrUnavailableError；
        引擎可用但全白页 → OCR 全空 → EmptyDocumentError。"""
        data = make_pdf_bytes([""])
        if _tesseract_available():
            with self.assertRaises(dx.EmptyDocumentError):
                dx.extract_text_from_bytes("scan.pdf", data)
        else:
            with self.assertRaises(dx.OcrUnavailableError):
                dx.extract_text_from_bytes("scan.pdf", data)


class TestOcrImage(unittest.TestCase):
    """图片 OCR（issue #50）：fitz 现场渲染文字图片（不依赖宿主字体），真实走 Tesseract。"""

    @unittest.skipUnless(_tesseract_available(), "本机无 tesseract 二进制")
    def test_eng_image(self):
        text = dx.extract_text_from_bytes("note.png", make_png_bytes("HELLO OCR 123"))
        self.assertIn("HELLO", text)
        self.assertIn("123", text)

    @unittest.skipUnless(_chi_sim_available(), "本机无 tesseract chi_sim 语言包")
    def test_chi_sim_image(self):
        text = dx.extract_text_from_bytes("采购.png", make_png_bytes("机密采购合同", fontname="china-s"))
        self.assertIn("机密", text)
        self.assertIn("合同", text)

    def test_tesseract_missing_degradation(self):
        """tesseract 二进制缺失 → OcrUnavailableError 明确中文报错（mock 执行点，不依赖真实引擎）。"""
        import pytesseract
        from unittest import mock

        with mock.patch("pytesseract.image_to_string",
                        side_effect=pytesseract.TesseractNotFoundError):
            with self.assertRaises(dx.OcrUnavailableError) as cm:
                dx.extract_text_from_bytes("note.png", make_png_bytes("HELLO"))
        self.assertIn("tesseract", str(cm.exception))

    def test_ocr_empty_result(self):
        """OCR 识别全空（低质量图）→ EmptyDocumentError（质量边界文案）。"""
        from unittest import mock

        with mock.patch("pytesseract.image_to_string", return_value="  \n "):
            with self.assertRaises(dx.EmptyDocumentError) as cm:
                dx.extract_text_from_bytes("note.png", make_png_bytes("HELLO"))
        # OCR 路径空结果给 OCR 质量提示（issue #51 P2-4 分路径文案的 OCR 侧）
        self.assertIn("OCR", str(cm.exception))

    def test_pixel_limit_image(self):
        """超大像素图片（issue #51 P1-1）：解码前按 Image.open 读到的 size 拦截（伪造 IHDR 尺寸，
        IDAT 仅 1 像素——拦截发生在 img.load() 前，无需真实大图），防超大图解码 OOM。"""
        data = make_png_with_dimensions(8000, 6000)  # 48M px > MAX_OCR_PIXELS
        with self.assertRaises(dx.OcrImageTooLargeError) as cm:
            dx.extract_text_from_bytes("big.png", data)
        self.assertIn("像素", str(cm.exception))
        self.assertIn("上限", str(cm.exception))

    def test_tesseract_timeout(self):
        """tesseract 超时（issue #51 P2-2）：pytesseract timeout kwarg 生效（断言传参），
        超时 RuntimeError → OcrUnavailableError 中文超时文案（mock 执行点，不依赖真实引擎）。"""
        from unittest import mock

        with mock.patch("pytesseract.image_to_string",
                        side_effect=RuntimeError("Tesseract process timeout")) as m:
            with self.assertRaises(dx.OcrUnavailableError) as cm:
                dx.extract_text_from_bytes("note.png", make_png_bytes("HELLO"))
        self.assertIn("超时", str(cm.exception))
        self.assertEqual(m.call_args.kwargs.get("timeout"), dx.OCR_PAGE_TIMEOUT)

    def test_image_internal_exception_wrapped(self):
        """图片 OCR 路径内部异常兜底（issue #51 P2-3）：pytesseract/Pillow 非引擎类异常
        （MemoryError/OSError/TypeError 等）→ CorruptDocumentError 中文文案，不裸泄。"""
        from unittest import mock

        with mock.patch("pytesseract.image_to_string", side_effect=MemoryError("boom")):
            with self.assertRaises(dx.CorruptDocumentError) as cm:
                dx.extract_text_from_bytes("note.png", make_png_bytes("HELLO"))
        self.assertIn("OCR", str(cm.exception))


class TestOcrScannedPdf(unittest.TestCase):
    """扫描 PDF OCR（issue #50）：内嵌文本为空时渲染页图逐页识别。"""

    @unittest.skipUnless(_tesseract_available(), "本机无 tesseract 二进制")
    def test_scanned_pdf_with_text_image(self):
        data = make_scanned_pdf_bytes(["CONFIDENTIAL SCAN 456"])
        text = dx.extract_text_from_bytes("scan.pdf", data)
        self.assertIn("CONFIDENTIAL", text)

    def test_page_limit(self):
        """OCR 页数上限 50：51 页无文本 PDF → 明确中文报错（在 OCR 调用前拦截，无需 tesseract）。"""
        data = make_pdf_bytes([""] * 51)
        with self.assertRaises(dx.EmptyDocumentError) as cm:
            dx.extract_text_from_bytes("long.pdf", data)
        self.assertIn("50", str(cm.exception))
        self.assertIn("页", str(cm.exception))

    def test_pixel_limit_huge_mediabox(self):
        """超大 MediaBox 畸形 PDF（issue #51 P1-1）：渲染前按 page.rect×DPI 系数拦截
        （在 get_pixmap 与 OCR 调用之前，无需 tesseract），且提取器不崩——后续正常提取照旧。"""
        data = make_huge_mediabox_pdf_bytes()  # 20000pt 见方 → 180DPI 渲染 25 亿 px
        with self.assertRaises(dx.OcrImageTooLargeError) as cm:
            dx.extract_text_from_bytes("evil.pdf", data)
        self.assertIn("像素", str(cm.exception))
        self.assertIn("上限", str(cm.exception))
        # 进程不崩：正常文档提取不受影响（对齐 #49 P1-1 超大文本用例的存活断言）
        text = dx.extract_text_from_bytes("ok.docx", make_docx_bytes(["正常文档内容仍然可用"]))
        self.assertIn("正常文档内容", text)


class TestOcrFallbackDensity(unittest.TestCase):
    """水印文本层密度启发式（issue #52 缺口 2）：有效字符/页 < 阈值且页面含扫描图才回退 OCR。"""

    @unittest.skipUnless(_chi_sim_available(), "本机无 tesseract chi_sim 语言包")
    def test_watermark_text_layer_triggers_ocr(self):
        """水印文本层 PDF（页图 + 每页「扫描全能王 创建」真文本水印）：非空文本层不挡回退，
        OCR 提取出正文（朴素「非空即返回」实现下指纹库只剩水印）。"""
        data = make_watermarked_scanned_pdf_bytes(["WATERMARK BODY EVIDENCE 789"])
        text = dx.extract_text_from_bytes("contract-scan.pdf", data)
        self.assertIn("WATERMARK", text)
        self.assertIn("789", text)

    def test_dense_digital_pdf_no_ocr(self):
        """密度 ≥ 阈值的数字 PDF（带内嵌图）不回退 OCR（不回归）：内嵌文本直接返回，
        pytesseract 全程不被调用（mock 守调用点，无需真实引擎）。"""
        from unittest import mock

        dense = "密度足够的数字文档正文内容" * 6  # >50 有效字符/页
        data = make_pdf_with_image_bytes(dense, fontname="china-s")
        with mock.patch("pytesseract.image_to_string",
                        side_effect=AssertionError("不应调用 OCR")):
            text = dx.extract_text_from_bytes("digital.pdf", data)
        self.assertIn("密度足够", text)

    def test_sparse_text_without_images_no_ocr(self):
        """稀疏文本但页面无扫描图 → 不回退 OCR（OCR 无可识别对象）：内嵌文本原样返回。"""
        from unittest import mock

        data = make_pdf_bytes(["短句"], fontname="china-s")  # 2 字符/页 < 阈值，但无图
        with mock.patch("pytesseract.image_to_string",
                        side_effect=AssertionError("不应调用 OCR")):
            text = dx.extract_text_from_bytes("sparse.pdf", data)
        self.assertIn("短句", text)

    def test_sparse_text_with_small_logo_no_ocr(self):
        """稀疏文本 + 小 logo 图不回退（issue #53 P1-1 误伤回归守）：单页 <50 字符 +
        面积仅 ~2% 的信头 logo/签名章小图——按面积比判定不是扫描页，内嵌文本原样返回
        （get_images() 式判定会误触发回退，把完美文本层换成 OCR 次品甚至 400 拒收）。"""
        from unittest import mock

        # make_pdf_with_image_bytes 内嵌 100×100pt 小图 ≈ 页面面积 2%（< SCAN_IMAGE_AREA_RATIO）
        data = make_pdf_with_image_bytes("报价单审批通过", fontname="china-s")
        with mock.patch("pytesseract.image_to_string",
                        side_effect=AssertionError("不应调用 OCR")):
            text = dx.extract_text_from_bytes("quote.pdf", data)
        self.assertIn("报价单审批通过", text)

    def test_mixed_density_average_boundary(self):
        """多页混合密度按全文平均判定（issue #53 P2-3 边界）：3 页（150/10/10 字符）平均
        56.7 ≥ 阈值 → 即便其中 2 页是整页扫描图也不回退（平均密度达标即信任内嵌文本层）。"""
        from unittest import mock

        data = make_mixed_density_pdf_bytes([
            ("密" * 150, False),               # 150 字符数字页
            ("稀疏十字符页一二三四五", True),    # 10 字符 + 整页扫描图
            ("稀疏十字符页六七八九十", True),    # 10 字符 + 整页扫描图
        ])
        with mock.patch("pytesseract.image_to_string",
                        side_effect=AssertionError("不应调用 OCR")):
            text = dx.extract_text_from_bytes("mixed.pdf", data)
        self.assertEqual(text.count("密"), 150)  # 盒式换行，按总数断言
        self.assertIn("稀疏十字符页一二三四五", text)


class TestEmptyResultMessage(unittest.TestCase):
    """空结果文案分路径（issue #51 P2-4）：非 OCR 空文档给原文案，不带 OCR 质量提示。"""

    def test_empty_docx_message_no_ocr_hint(self):
        with self.assertRaises(dx.EmptyDocumentError) as cm:
            dx.extract_text_from_bytes("blank.docx", make_docx_bytes(["", "  "]))
        self.assertIn("未提取到文本", str(cm.exception))
        self.assertNotIn("OCR", str(cm.exception))

    def test_empty_txt_message_no_ocr_hint(self):
        with self.assertRaises(dx.EmptyDocumentError) as cm:
            dx.extract_text_from_bytes("blank.txt", b"   \n  ")
        self.assertNotIn("OCR", str(cm.exception))


class TestExtractedTextCap(unittest.TestCase):
    """提取文本 8M 字符上限（issue #49 P1-1）：zip/流扩张防 OOM，超限明确中文报错。"""

    def test_over_cap_rejected(self):
        data = ("x" * (dx.MAX_EXTRACTED_CHARS + 1)).encode("utf-8")  # ~8MB，线路 16MB 上限内
        with self.assertRaises(dx.ExtractedTextTooLargeError) as cm:
            dx.extract_text_from_bytes("big.txt", data)
        self.assertIn("上限", str(cm.exception))

    def test_just_under_cap_passes(self):
        line = "商密文档边界行 abcdefghijklmnopqrstuvwxyz0123456789"  # >12 字符，避免误判场景干扰
        data = (line * 100).encode("utf-8")  # 远小于上限的正常文档
        text = dx.extract_text_from_bytes("ok.txt", data)
        self.assertIn("商密文档边界行", text)


if __name__ == "__main__":
    unittest.main()
